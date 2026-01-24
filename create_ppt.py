#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
포트폴리오 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 색상 정의
KAKAO_YELLOW = RGBColor(254, 229, 0)
KAKAO_BLACK = RGBColor(25, 25, 25)
TEXT_SUB = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """표지 슬라이드"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경색 (카카오 옐로우)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = KAKAO_YELLOW
    background.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SUB
    p.alignment = PP_ALIGN.CENTER

def add_about_slide(prs):
    """About 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 섹션 라벨
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ABOUT"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_SUB
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "뉴미디어 채널과 콘텐츠로\n브랜드를 성장시키는 IMC 마케터"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    
    # 설명
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "라이프스타일 패션 브랜드에서 IMC 마케터로 일하며 브랜딩과 매출이라는 두 축을 실질적인 성과로 연결해왔습니다. 데이터와 감각을 균형감 있게 다루며 팀과 속도를 맞춰가는 방식으로, 브랜드의 성장은 팀 전체가 목표를 함께 향할 때 가능하다고 믿습니다."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_SUB
    
    # Stats
    stats = [
        ("4+", "Years at F&F"),
        ("500만", "Max Contents View"),
        ("손익분기점 달성", "신규 브랜드")
    ]
    
    for i, (num, label) in enumerate(stats):
        x = Inches(0.5 + i * 3)
        
        # 숫자
        num_box = slide.shapes.add_textbox(x, Inches(4), Inches(2.5), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = KAKAO_BLACK
        
        # 라벨
        label_box = slide.shapes.add_textbox(x, Inches(4.7), Inches(2.5), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_SUB

def add_capabilities_slide(prs):
    """핵심 역량 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 섹션 라벨
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPABILITIES"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_SUB
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "핵심 역량"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    
    capabilities = [
        ("01", "IMC 전략 & 브랜드 성장", "SNS·PR·셀럽·팝업·광고를 하나의 메시지로 통합\n6개월 만에 인지도·매출 2배 성장"),
        ("02", "데이터 기반 퍼포먼스", "CTR/CVR/ROAS 분석, AI 툴 활용\n광고비 20% 절감 + 전환율 상승"),
        ("03", "온·오프라인 리테일", "더현대·EQL·롯데월드몰 팝업 A to Z\n검색 유입 150~400% 증가"),
        ("04", "AI 기반 업무 자동화", "SQL 기반 KPI 대시보드 구축\n회의 자료 제작 시간 3배 단축")
    ]
    
    for i, (num, title, desc) in enumerate(capabilities):
        row = i // 2
        col = i % 2
        x = Inches(0.5 + col * 4.8)
        y = Inches(1.6 + row * 2)
        
        # 번호
        num_box = slide.shapes.add_textbox(x, y, Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = KAKAO_BLACK
        
        # 제목
        title_box = slide.shapes.add_textbox(x + Inches(0.5), y, Inches(4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = KAKAO_BLACK
        
        # 설명
        desc_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(0.4), Inches(4), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_SUB

def add_project_slide(prs, index, name, category, period, summary, role, problem, solution, results, insight):
    """프로젝트 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 인덱스
    idx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(0.6), Inches(0.4))
    tf = idx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    
    # 카테고리 & 기간
    meta_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.3), Inches(6), Inches(0.4))
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{category}  |  {period}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_SUB
    
    # 프로젝트명
    name_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.6))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    
    # 요약
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.8))
    tf = summary_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = summary
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_SUB
    
    # 좌측: Role + Problem
    left_x = Inches(0.5)
    
    # My Role
    role_title = slide.shapes.add_textbox(left_x, Inches(2.2), Inches(4.3), Inches(0.3))
    tf = role_title.text_frame
    p = tf.paragraphs[0]
    p.text = "👤 My Role"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    
    role_box = slide.shapes.add_textbox(left_x, Inches(2.5), Inches(4.3), Inches(1.2))
    tf = role_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SUB
    
    # Problem
    prob_title = slide.shapes.add_textbox(left_x, Inches(3.7), Inches(4.3), Inches(0.3))
    tf = prob_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔍 Problem"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(225, 29, 72)
    
    prob_box = slide.shapes.add_textbox(left_x, Inches(4), Inches(4.3), Inches(0.9))
    tf = prob_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = problem
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SUB
    
    # 우측: Solution + Result
    right_x = Inches(5.2)
    
    # Solution
    sol_title = slide.shapes.add_textbox(right_x, Inches(2.2), Inches(4.3), Inches(0.3))
    tf = sol_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 Solution"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(37, 99, 235)
    
    sol_box = slide.shapes.add_textbox(right_x, Inches(2.5), Inches(4.3), Inches(1.2))
    tf = sol_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = solution
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SUB
    
    # Result
    res_title = slide.shapes.add_textbox(right_x, Inches(3.7), Inches(4.3), Inches(0.3))
    tf = res_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📈 Result"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 163, 74)
    
    res_box = slide.shapes.add_textbox(right_x, Inches(4), Inches(4.3), Inches(0.9))
    tf = res_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = results
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SUB
    
    # Insight (하단)
    ins_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(9), Inches(0.3))
    tf = ins_title.text_frame
    p = tf.paragraphs[0]
    p.text = "💭 Insight"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(217, 119, 6)
    
    ins_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(9), Inches(0.5))
    tf = ins_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = insight
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_SUB

def add_contact_slide(prs):
    """Contact 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 섹션 라벨
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTACT"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_SUB
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Let's Work\nTogether"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = KAKAO_BLACK
    
    # 설명
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(7), Inches(1.2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "뉴미디어 채널과 콘텐츠로 브랜드를 성장시키는 IMC 마케터입니다.\n브랜드 커뮤니케이션 전략 수립부터 통합 마케팅, 국내외 협업 캠페인까지\n채널 × 콘텐츠 × 데이터를 연결해 브랜드를 빌드업합니다."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_SUB

def main():
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 비율
    
    # 1. 표지
    add_title_slide(prs, "Brand Maker : 이난주", "CONTENTS · DATA-BASED STRATEGY · AI\nF&F · 브랜드 마케팅 · 4년+")
    
    # 2. About
    add_about_slide(prs)
    
    # 3. 핵심 역량
    add_capabilities_slide(prs)
    
    # 4. 프로젝트들
    projects = [
        {
            "index": 1,
            "name": "듀베티카 24FW 셀럽/인플루언서 협업 프로젝트",
            "category": "셀럽/인플루언서 · IMC · PM",
            "period": "2024.09 - 2025.02",
            "summary": "[PM] 리브랜딩 이후 듀베티카가 가진 패딩 헤리티지/프리미엄 무드를 대중에게 확산시키고, 'NEXT 몽클레르' 포지셔닝을 만들기 위해 셀럽(김지원) → 인플루언서 → 커뮤니티 → 실수요층으로 이어지는 IMC 메시지 통합을 설계·실행한 프로젝트",
            "role": "• 기획 (80%): 타깃 재정의, 포지셔닝, 메시지 통합 설계\n• 실행 (100%): 셀럽 활용, 팝업, 옥외광고, 바이럴 연결\n• 분석 (100%): 노출→인스타→커뮤니티→실수요 퍼널 추적",
            "problem": "• 브랜드 인지도 개선 필요\n• 아웃도어 치우친 톤앤매너 재정립\n• 마케팅 타깃 연령 30대 여성으로 낮추기",
            "solution": "• 셀럽 앵커: 김지원 이미지로 포지셔닝 직관화\n• 연결고리 설계: 셀럽→팝업→인플루언서 체인\n• 키 아이템 메이킹: 판쵸 스타일 중심 선택/설득",
            "results": "• 셀럽 콘텐츠 조회수 500만\n• 키 아이템 매출 4배 증가\n• 연간 검색량 & 매출 2배",
            "insight": "디테일이 브랜드를 만든다 — 모든 터치포인트를 하나의 문장처럼 연결할 때 전달력이 커진다"
        },
        {
            "index": 2,
            "name": "세르지오 타키니 한남 플래그십 & 콜라보",
            "category": "Flagship · Collaboration · PM",
            "period": "2025.08 - 2025.10",
            "summary": "[PM 총괄] 60년 헤리티지를 가진 스포츠 브랜드를 클래식 웰니스 라이프스타일 브랜드로 확장. 앰버서더 박지현 이미지와 한남 플래그십을 거점으로 2030 여성 타깃 유입 및 매출 전환 2배 달성",
            "role": "• 기획 리드 (100%): 플래그십 오픈, 콜라보 기획\n• 실행 (85%): 현장 운영, 셀럽/인플루언서 마케팅\n• 분석 (100%): 동선/체험/촬영 스팟 설계, 성과 분석",
            "problem": "• 테니스 브랜드로 고착된 이미지 개선\n• 헤리티지 중심 톤앤매너 재정립\n• 마케팅 타깃 20대 여성으로 전환",
            "solution": "• 앰버서더 박지현으로 브랜드 리프레시\n• 한남 플래그십을 확산 거점화\n• 라이프스타일 카테고리(셋업/니트) 확장",
            "results": "• 한남 플래그십 유입·매출 전환 2배\n• 단일 콘텐츠 조회수 300만+\n• 키 아이템 월간 매출 TOP 5\n• 오가닉 콘텐츠 약 400건 생성",
            "insight": "브랜드 강/약점 진단 → 타깃 커뮤니티가 납득할 가치 제시 → 내외부 메시지 일관성이 성과를 좌우"
        },
        {
            "index": 3,
            "name": "현대카드 슈퍼매치",
            "category": "[ 카테고리 ]",
            "period": "[ 기간 ]",
            "summary": "[ 프로젝트 요약 - 내용 입력 필요 ]",
            "role": "[ 역할 입력 필요 ]",
            "problem": "[ 문제 입력 필요 ]",
            "solution": "[ 솔루션 입력 필요 ]",
            "results": "[ 결과 입력 필요 ]",
            "insight": "[ 인사이트 입력 필요 ]"
        },
        {
            "index": 4,
            "name": "세르지오 타키니 EQL 성수 팝업",
            "category": "Retail Pop-up · PM",
            "period": "2025.02 - 2025.04",
            "summary": "[PM 리드] 성수 상권에 맞는 감도로 브랜드 재해석. 타겟 연령층 2030으로 확장하고 SNS 노출 극대화를 위한 체험형 공간을 기획부터 현장 운영까지 전 과정 주도",
            "role": "• 기획 리드 (100%): 팝업 비주얼, 브랜드 재해석\n• 실행 (90%): 오픈행사 운영, 인플루언서 초청\n• 분석 (100%): 고객 동선, 촬영 스팟 큐레이션",
            "problem": "• 성수 상권에 맞는 브랜드 재해석 필요\n• 4050 → 2030 연령층 확장 필요\n• 빠른 트렌드 속 명확한 첫인상 필요",
            "solution": "• '고객의 촬영 행동' 중심 공간 설계\n• 제품 2~3개로 최소화, 무드 일관성\n• 인플루언서 초청 + SNS 바이럴 기획",
            "results": "• SNS 바이럴 조회수 300만+\n• 키 아이템 검색량 2배\n• 글로벌 고객층 확장",
            "insight": "공간의 역할 재정의 — 성수에서는 오프라인 목적이 판매보다 '온라인 확산 가능한 브랜드 경험'이어야 함"
        },
        {
            "index": 5,
            "name": "몬테카를로 마스터즈 MCCC 글로벌 행사",
            "category": "Global Event · PM",
            "period": "2025.02 - 2025.05",
            "summary": "[PM 리드] 글로벌 테니스 대회 스폰서십을 활용한 브랜드 테니스 헤리티지 강화 프로젝트. 셀럽 초청 이벤트 기획 및 멀티 유즈 관점의 콘텐츠 자산 확보 주도",
            "role": "• 기획 리드 (100%): 셀럽 초청 타임라인, 동선/착장/촬영 설계\n• 실행 (70%): 현장 운영, 해외 행사 운영비 관리\n• 분석 (100%): PR·바이럴 전략, 채널별 촬영 구도 설계",
            "problem": "• 브랜드 테니스 헤리티지 강화 필요\n• 해외/국내 팬층 내 인지도 확장\n• 제약조건 속 글로벌 콘텐츠 자산 확보",
            "solution": "• 셀럽 초청 행사 타임라인 및 현장 운영 기획\n• 온/오프라인 PR, 바이럴 전략 수립\n• 채널별 멀티 유즈 촬영 구도 설계",
            "results": "• 메인 콘텐츠 ENG 3.5만\n• 주요 콘텐츠 조회수 100~300만\n• 고퀄리티 콘텐츠 자산 확보",
            "insight": "멀티 유즈 관점의 기획 — 현장의 화려함보다 '얼마나 오래, 여러 채널에서 쓸 수 있는지'가 중요"
        },
        {
            "index": 6,
            "name": "시즌 IMC 전략 수립",
            "category": "Strategy · PM",
            "period": "2021.10 - Present",
            "summary": "[전략 리드] 브랜딩과 세일즈가 동시에 작동하는 통합 마케팅 전략을 시즌별로 리드. 데이터 기반 KPI 설정과 예산 배분으로 지속 가능한 성장 구조 구축",
            "role": "• 전략 리드 (100%): 시즌별 상품 전략, KPI, 예산, 플래닝\n• 데이터 분석 (100%): 트렌드 리서치 자료 배포\n• 조직 조율 (60%): CEO 보고, 부서 회의체 운영",
            "problem": "• 채널, 예산, 메시지가 부서별로 분산\n• 감에 의존한 의사결정으로 우선순위 흔들림\n• 조직 내 공통 참고 지표 부재",
            "solution": "• 시즌별 상품 전략, KPI, 예산, 플래닝 리드\n• 데이터 기반 전략 및 트렌드 리서치 배포\n• MD/영업/온라인 등 세일즈 부서 협업 주도",
            "results": "• 시즌별 인지도·매출 2배+\n• 전략/리서치/KPI 자료 템플릿화\n• 리드 타임 단축",
            "insight": "구조의 힘 — '좋은 아이디어'보다 '지속 가능한 성장 구조'가 더 중요. 조직을 같은 방향으로 움직이게 하는 브릿지 역할"
        },
        {
            "index": 7,
            "name": "AI Agent MKT KPI 대시보드 구축",
            "category": "AI · Automation",
            "period": "2025.08 - Present",
            "summary": "[개발 리드] 마케팅 KPI와 판매데이터를 연결하는 AI 기반 대시보드를 SQL 쿼리 설계부터 대시보드 제작까지 직접 개발. 회의 자료 제작 시간 3배 단축",
            "role": "• 기획 리드 (100%): MKT KPI-판매데이터 연결 구조 설계\n• 개발 (100%): Snowflake SQL 쿼리, AI 대시보드 제작\n• 자동화 구축 (100%): 주간 회의, CEO 보고 자료 자동화",
            "problem": "• 마케팅 KPI와 판매데이터 연결 부재\n• 주간 회의, 보고 문서 제작에 과도한 시간\n• 실시간 KPI 기반 의사결정 환경 미비",
            "solution": "• Snowflake SQL 데이터 자동화 쿼리 설계\n• AI 기반 인사이트 도출 KPI 대시보드 제작\n• 경쟁사 대비 SNS, 판매, 재고 데이터 통합",
            "results": "• 자료 제작 시간 3배 단축\n• 의사결정 속도 향상\n• 부서 커뮤니케이션 감소",
            "insight": "AI는 팀의 운영 전략 — AI는 단순히 도구가 아니라 팀의 일하는 방식을 바꾸는 운영 전략"
        }
    ]
    
    for proj in projects:
        add_project_slide(prs, **proj)
    
    # 마지막: Contact
    add_contact_slide(prs)
    
    # 저장
    output_path = "/Users/leenanju/portpolio/이난주_포트폴리오.pptx"
    prs.save(output_path)
    print(f"✅ PPT 파일 생성 완료: {output_path}")

if __name__ == "__main__":
    main()
